from __future__ import annotations

import sys
from pathlib import Path

ROOT = Path(r"D:\Documents\Workspace\traceability-system")
sys.path.insert(0, str(ROOT / ".pptx_deps"))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from PIL import Image


ASSETS = ROOT / "slide_assets"
ICONS = ASSETS / "icons_v2"
OUT = ROOT / "DATN_Nguyen_Dang_Nam_Traceability_Blockchain_clean.pptx"


COL = {
    "bg": "F7FAFC",
    "paper": "FFFFFF",
    "ink": "0B1220",
    "text": "1F2937",
    "muted": "64748B",
    "line": "D8E0EA",
    "blue": "2563EB",
    "blue_soft": "EAF2FF",
    "green": "16A34A",
    "green_soft": "EAFBF0",
    "orange": "F97316",
    "orange_soft": "FFF3E8",
    "red": "DC2626",
    "red_soft": "FFF0F0",
    "dark": "0F172A",
}


def rgb(hex_color: str) -> RGBColor:
    h = hex_color.strip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def set_fill(shape, color: str):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)


def set_line(shape, color: str | None = None, width: float = 1.0):
    if color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(color)
        shape.line.width = Pt(width)


def add_box(slide, x, y, w, h, fill="FFFFFF", line="D8E0EA", radius=True):
    shp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    set_fill(shp, fill)
    set_line(shp, line)
    return shp


def add_text(slide, x, y, w, h, text, size=18, color="1F2937", bold=False, align="left", valign="top"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    tf.vertical_anchor = {
        "top": MSO_ANCHOR.TOP,
        "mid": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }.get(valign, MSO_ANCHOR.TOP)
    lines = str(text).split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
        }.get(align, PP_ALIGN.LEFT)
        for run in p.runs:
            run.font.name = "Aptos"
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = rgb(color)
    return tb


def add_icon(slide, name: str, x, y, size, bg="2563EB"):
    circ = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
    set_fill(circ, bg)
    set_line(circ, None)
    p = ICONS / f"{name}.png"
    if p.exists():
        pad = size * 0.22
        slide.shapes.add_picture(str(p), Inches(x + pad), Inches(y + pad), Inches(size - 2 * pad), Inches(size - 2 * pad))
    return circ


def add_card(slide, x, y, w, h, title, body, accent="2563EB", bg="FFFFFF", icon_name=None):
    add_box(slide, x, y, w, h, bg, "E1E8F0")
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    set_fill(bar, accent)
    set_line(bar, None)
    tx = x + 0.28
    tw = w - 0.48
    if icon_name:
        add_icon(slide, icon_name, x + 0.28, y + 0.24, 0.52, accent)
        tx = x + 0.95
        tw = w - 1.18
    add_text(slide, tx, y + 0.20, tw, 0.32, title, 15, COL["ink"], True)
    add_text(slide, tx, y + 0.62, tw, h - 0.72, body, 12, COL["muted"])


def add_header(slide, title, subtitle, number):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(COL["bg"])
    left = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(0.13), Inches(7.5))
    set_fill(left, COL["blue"])
    set_line(left, None)
    add_text(slide, 0.55, 0.28, 9.8, 0.4, title, 25, COL["ink"], True)
    add_text(slide, 0.57, 0.78, 10.8, 0.28, subtitle, 12, COL["muted"])
    add_text(slide, 12.35, 0.35, 0.45, 0.28, f"{number:02d}", 11, COL["muted"], True, "right")
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(1.13), Inches(1.0), Inches(0.045))
    set_fill(line, COL["green"])
    set_line(line, None)


def add_bullets(slide, x, y, items, color="16A34A", size=14, gap=0.48):
    yy = y
    for item in items:
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(yy + 0.12), Inches(0.12), Inches(0.12))
        set_fill(dot, color)
        set_line(dot, None)
        add_text(slide, x + 0.28, yy, 5.9, 0.38, item, size, COL["text"])
        yy += gap


def image_path(name: str) -> Path | None:
    p = ASSETS / name
    return p if p.exists() else None


def add_fit_picture(slide, path: Path, x, y, w, h):
    with Image.open(path) as img:
        iw, ih = img.size
    scale = min(w / iw, h / ih)
    nw = iw * scale
    nh = ih * scale
    slide.shapes.add_picture(str(path), Inches(x + (w - nw) / 2), Inches(y + (h - nh) / 2), Inches(nw), Inches(nh))


def phone_mock(slide, path: Path, x, y, w, h):
    add_box(slide, x, y, w, h, COL["dark"], None)
    add_fit_picture(slide, path, x + 0.12, y + 0.22, w - 0.24, h - 0.44)


def chip(slide, x, y, label, color):
    add_box(slide, x, y, 1.45, 0.33, color, None)
    add_text(slide, x + 0.08, y + 0.055, 1.29, 0.16, label, 9, "FFFFFF", True, "center")


def make_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = rgb(COL["bg"])
    top = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333333), Inches(0.16))
    set_fill(top, COL["blue"])
    set_line(top, None)
    add_box(s, 0.6, 0.66, 6.95, 5.85, COL["paper"], "E1E8F0")
    add_text(s, 1.02, 1.05, 6.1, 1.65, "Xây dựng hệ thống truy xuất nguồn gốc hàng hóa\nsử dụng công nghệ blockchain và microservice", 28, COL["ink"], True)
    add_text(s, 1.04, 2.98, 5.7, 0.34, "Đồ án tốt nghiệp ngành Công nghệ thông tin", 15, COL["blue"], True)
    bar = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.04), Inches(3.48), Inches(1.35), Inches(0.055))
    set_fill(bar, COL["green"])
    set_line(bar, None)
    add_text(s, 1.04, 4.0, 5.9, 0.62, "Nguyễn Đăng Nam  |  CT060226  |  CT6B", 16, COL["text"], True)
    add_text(s, 1.04, 4.72, 5.7, 0.7, "GVHD: TS. Trần Anh Tú\nThS. Nguyễn Văn Quyết", 14, COL["muted"])
    phone = image_path("image54.png") or image_path("image49.png")
    if phone:
        phone_mock(s, phone, 8.25, 0.72, 2.65, 5.9)
    qr = image_path("image37.png")
    if qr:
        add_box(s, 10.98, 3.9, 1.65, 1.2, COL["paper"], "E1E8F0")
        add_fit_picture(s, qr, 11.13, 4.07, 1.35, 0.75)
    chip(s, 8.15, 6.62, "Blockchain", COL["blue"])
    chip(s, 9.78, 6.62, "Microservice", COL["green"])
    chip(s, 11.4, 6.62, "QR/Serial", COL["orange"])

    # Slide 2
    s = prs.slides.add_slide(blank)
    add_header(s, "Bài toán đặt ra", "Tại điểm mua, người dùng cần thông tin đáng tin và dễ kiểm chứng", 2)
    add_card(s, 0.75, 1.55, 3.8, 2.1, "Khó kiểm chứng", "Người mua khó biết sản phẩm thật sự đến từ đâu và đã đi qua những khâu nào.", COL["blue"], COL["paper"], "warning")
    add_card(s, 4.82, 1.55, 3.8, 2.1, "Mã có thể bị sao chép", "Thông tin in trên bao bì hoặc mã truy xuất thông thường có thể bị dùng lại.", COL["orange"], COL["paper"], "scan")
    add_card(s, 8.88, 1.55, 3.8, 2.1, "Dữ liệu cần kiểm chứng", "Nếu dữ liệu trong hệ thống bị sửa, cần có cách phát hiện sai lệch.", COL["green"], COL["paper"], "shield")
    add_box(s, 1.1, 4.65, 11.1, 1.18, COL["blue_soft"], "C9DDFE")
    add_text(s, 1.45, 4.95, 10.4, 0.5, "Hướng giải quyết: xây dựng hệ thống truy xuất nguồn gốc có mobile scan, dữ liệu chuỗi cung ứng và blockchain lưu hash để xác minh.", 19, COL["ink"], True, "center", "mid")

    # Slide 3
    s = prs.slides.add_slide(blank)
    add_header(s, "Mục tiêu và ý nghĩa thực tiễn", "Tập trung vào trải nghiệm người mua hàng", 3)
    add_card(s, 0.78, 1.55, 3.82, 2.35, "Xem thông tin tại siêu thị", "Quét QR hoặc nhập serial để xem nguồn gốc, lịch sử vận chuyển và chi tiết sản phẩm.", COL["blue"], COL["paper"], "mobile")
    add_card(s, 4.82, 1.55, 3.82, 2.35, "Tham khảo đánh giá", "Người mua sau có thể xem đánh giá từ người đã mua và sử dụng sản phẩm trước đó.", COL["green"], COL["paper"], "review")
    add_card(s, 8.88, 1.55, 3.82, 2.35, "Hỗ trợ chống giả", "Mỗi sản phẩm gắn với một serial tờ tiền đã đăng ký; serial chỉ dùng một lần để phát hiện dùng lại hoặc không khớp.", COL["orange"], COL["paper"], "shield")
    add_text(s, 1.0, 4.65, 11.4, 0.92, "Mục tiêu kỹ thuật: quản lý chuỗi cung ứng nhiều vai trò, truy xuất ngược đến nguyên liệu và xác minh tính toàn vẹn dữ liệu bằng blockchain.", 21, COL["ink"], True, "center", "mid")

    # Slide 4
    s = prs.slides.add_slide(blank)
    add_header(s, "Phạm vi và tác nhân", "Mỗi vai trò tương ứng một lát cắt nghiệp vụ trong hệ thống", 4)
    roles = [
        ("Admin", "Duyệt vai trò,\nquản lý người dùng", COL["blue"], "admin"),
        ("Supplier", "Tạo và quản lý\nlô nguyên liệu", COL["green"], "leaf"),
        ("Manufacturer", "Sản xuất, đóng gói,\nđăng ký serial", COL["orange"], "factory"),
        ("Retailer", "Đặt hàng\nthành phẩm", COL["blue"], "retail"),
        ("Transporter", "Xác nhận lấy hàng\nvà giao hàng", COL["green"], "truck"),
        ("User", "Quét truy xuất,\nxem/đánh giá SP", COL["orange"], "user"),
    ]
    for i, (name, body, accent, ic) in enumerate(roles):
        add_card(s, 0.78 + (i % 3) * 4.15, 1.55 + (i // 3) * 2.15, 3.72, 1.65, name, body, accent, COL["paper"], ic)

    # Slide 5
    s = prs.slides.add_slide(blank)
    add_header(s, "Kiến trúc tổng thể", "Thiết kế theo hướng microservices, có gateway và xử lý bất đồng bộ", 5)
    add_card(s, 0.8, 1.65, 2.6, 1.2, "React Web", "Dashboard nghiệp vụ", COL["blue"], COL["paper"], "web")
    add_card(s, 0.8, 3.35, 2.6, 1.2, "Flutter Mobile", "Scan, retailer, transporter", COL["green"], COL["paper"], "mobile")
    for yy in [2.28, 3.98]:
        line = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(3.55), Inches(yy), Inches(0.75), Inches(0.06))
        set_fill(line, COL["line"])
        set_line(line, None)
    add_card(s, 4.35, 2.42, 2.3, 1.18, "API Gateway", "Điểm vào chung", COL["orange"], COL["orange_soft"], "server")
    line = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(6.78), Inches(3.0), Inches(0.65), Inches(0.06))
    set_fill(line, COL["line"])
    set_line(line, None)
    for i, svc in enumerate(["Identity", "Catalog", "Core", "Trade", "Blockchain"]):
        x = 7.55 + (i % 3) * 1.58
        y = 1.55 + (i // 3) * 1.52
        add_box(s, x, y, 1.35, 0.76, COL["paper"], "D7E2EF")
        add_text(s, x + 0.05, y + 0.26, 1.25, 0.2, svc, 11, COL["ink"], True, "center")
    add_card(s, 7.55, 4.65, 1.85, 1.25, "PostgreSQL", "Lưu nghiệp vụ", COL["blue"], COL["blue_soft"], "db")
    add_card(s, 9.72, 4.65, 1.85, 1.25, "Kafka", "Event async", COL["green"], COL["green_soft"], "kafka")
    add_card(s, 11.88, 4.65, 1.1, 1.25, "Chain", "Hash", COL["orange"], COL["orange_soft"], "chain")

    # Slide 6
    s = prs.slides.add_slide(blank)
    add_header(s, "Công nghệ sử dụng", "Chia theo các lớp triển khai chính", 6)
    stacks = [
        ("Backend", "Spring Boot\nSpring Security\nEureka, Feign", COL["blue"], COL["blue_soft"], "server"),
        ("Frontend Web", "React, Vite\nAnt Design\nReact Query", COL["green"], COL["green_soft"], "web"),
        ("Mobile", "Flutter, Bloc\nDio\nQR/OCR Scanner", COL["orange"], COL["orange_soft"], "mobile"),
        ("Blockchain & Hạ tầng", "Solidity, Web3j\nPostgreSQL, Kafka\nWebSocket, Ganache", COL["blue"], COL["blue_soft"], "chain"),
    ]
    for i, (title, body, accent, bg, ic) in enumerate(stacks):
        add_card(s, 0.78 + i * 3.05, 1.75, 2.72, 3.35, title, body, accent, bg, ic)
    add_text(s, 1.05, 5.9, 11.2, 0.48, "Lựa chọn công nghệ phục vụ ba yêu cầu: chia tách nghiệp vụ, truy xuất trên thiết bị di động và xác minh dữ liệu bằng blockchain.", 17, COL["ink"], True, "center")

    # Slide 7
    s = prs.slides.add_slide(blank)
    add_header(s, "Luồng nghiệp vụ chính", "Tóm tắt đường đi của hàng hóa trong hệ thống", 7)
    steps = [
        ("01", "Supplier\ntạo nguyên liệu", COL["green"]),
        ("02", "Manufacturer\nmua nguyên liệu", COL["blue"]),
        ("03", "Sản xuất\npallet", COL["orange"]),
        ("04", "Đóng gói\ncarton/unit", COL["blue"]),
        ("05", "Retailer\nđặt hàng", COL["green"]),
        ("06", "Transporter\ngiao hàng", COL["orange"]),
        ("07", "User\nquét truy xuất", COL["blue"]),
    ]
    for i, (num, label, color) in enumerate(steps):
        x = 0.58 + i * 1.78
        circ = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(2.0), Inches(0.86), Inches(0.86))
        set_fill(circ, color)
        set_line(circ, None)
        add_text(s, x, 2.28, 0.86, 0.2, num, 12, "FFFFFF", True, "center")
        add_text(s, x - 0.38, 3.15, 1.58, 0.65, label, 13, COL["ink"], True, "center")
        if i < len(steps) - 1:
            l = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + 0.95), Inches(2.42), Inches(0.68), Inches(0.055))
            set_fill(l, COL["line"])
            set_line(l, None)
    add_box(s, 0.92, 5.05, 11.55, 1.0, COL["paper"], "E1E8F0")
    add_text(s, 1.25, 5.34, 10.9, 0.42, "Điểm đáng chú ý: hệ thống quản lý cả dữ liệu truy xuất và trạng thái giao dịch/vận chuyển giữa nhiều vai trò.", 18, COL["text"], True, "center")

    # Slide 8
    s = prs.slides.add_slide(blank)
    add_header(s, "Mô hình truy xuất nguồn gốc", "Truy ngược từ một sản phẩm lẻ về nguyên liệu ban đầu", 8)
    nodes = [
        ("Raw Batch", "Lô nguyên liệu", COL["green"], COL["green_soft"]),
        ("Pallet", "Lô sản xuất", COL["green"], COL["green_soft"]),
        ("Carton", "Thùng hàng", COL["blue"], COL["blue_soft"]),
        ("Product Unit", "Sản phẩm lẻ", COL["blue"], COL["blue_soft"]),
        ("QR / Serial", "Mã truy xuất", COL["orange"], COL["orange_soft"]),
        ("User Scan", "Người dùng", COL["orange"], COL["orange_soft"]),
    ]
    for i, (title, sub, col, bg) in enumerate(nodes):
        x = 0.6 + i * 2.08
        add_box(s, x, 2.0, 1.5, 0.96, bg, "DCE5EF")
        add_text(s, x + 0.08, 2.22, 1.34, 0.2, title, 12, COL["ink"], True, "center")
        add_text(s, x + 0.08, 2.55, 1.34, 0.18, sub, 9, COL["muted"], False, "center")
        if i < 5:
            l = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + 1.58), Inches(2.45), Inches(0.38), Inches(0.055))
            set_fill(l, col)
            set_line(l, None)
    add_card(s, 0.95, 4.35, 3.7, 1.35, "Dữ liệu nghiệp vụ", "Product, carton, pallet, raw batch và transfer record tạo thành timeline.", COL["blue"], COL["paper"], "db")
    add_card(s, 4.95, 4.35, 3.7, 1.35, "Blockchain", "Hash dùng để kiểm chứng tính toàn vẹn của dữ liệu.", COL["green"], COL["paper"], "chain")
    add_card(s, 8.95, 4.35, 3.7, 1.35, "Mobile scan", "Người dùng xem timeline và kết quả xác minh.", COL["orange"], COL["paper"], "scan")

    # Slide 9
    s = prs.slides.add_slide(blank)
    add_header(s, "Cơ chế blockchain và xác minh", "Không đưa toàn bộ dữ liệu lên blockchain, chỉ neo hash quan trọng", 9)
    flow = [
        ("Dữ liệu nghiệp vụ", "Lưu trong PostgreSQL", "db", COL["blue"]),
        ("Tính hash", "Chuẩn hóa và băm dữ liệu", "shield", COL["green"]),
        ("Ghi blockchain", "Lưu hash + event audit", "chain", COL["orange"]),
        ("Verify", "Tính lại và so sánh", "shield", COL["blue"]),
    ]
    for i, (title, body, ic, col) in enumerate(flow):
        x = 0.82 + i * 3.08
        add_card(s, x, 1.85, 2.55, 1.5, title, body, col, COL["paper"], ic)
        if i < 3:
            l = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + 2.66), Inches(2.58), Inches(0.44), Inches(0.055))
            set_fill(l, COL["line"])
            set_line(l, None)
    add_card(s, 1.08, 4.35, 5.35, 1.36, "Khi dữ liệu còn nguyên vẹn", "Hash tính lại từ DB khớp với hash đã ghi trên blockchain.", COL["green"], COL["green_soft"], "shield")
    add_card(s, 6.88, 4.35, 5.35, 1.36, "Khi dữ liệu bị sửa", "Hash không khớp, hệ thống cảnh báo dữ liệu truy xuất không còn toàn vẹn.", COL["red"], COL["red_soft"], "warning")

    # Slide 10
    s = prs.slides.add_slide(blank)
    add_header(s, "Chống giả và đánh giá sản phẩm", "Hai chức năng giúp người dùng có thêm niềm tin khi mua hàng", 10)
    add_card(s, 0.78, 1.55, 5.55, 2.05, "Chống giả bằng serial tờ tiền", "Mỗi sản phẩm gắn với một serial tờ tiền đã đăng ký; serial chỉ được dùng một lần.\nNếu serial bị dùng lại hoặc không khớp sản phẩm, hệ thống phát hiện bất thường.", COL["orange"], COL["orange_soft"], "shield")
    add_card(s, 0.78, 3.95, 5.55, 1.65, "Đánh giá sản phẩm", "Người dùng đã mua có thể đánh giá, người mua sau có thêm thông tin tham khảo thực tế.", COL["green"], COL["green_soft"], "review")
    demo1 = image_path("image58.png")
    demo2 = image_path("image53.png") or image_path("image49.png")
    if demo1:
        phone_mock(s, demo1, 7.2, 1.42, 1.82, 4.15)
    if demo2:
        phone_mock(s, demo2, 9.45, 1.42, 1.82, 4.15)

    # Slide 11
    s = prs.slides.add_slide(blank)
    add_header(s, "Kết quả đạt được", "Các chức năng chính đã hoàn thiện để demo bảo vệ", 11)
    pics = [
        ("image41.png", 0.72, 1.45, 3.7, 1.24),
        ("image40.png", 4.65, 1.45, 3.7, 1.75),
        ("image62.png", 8.72, 1.28, 1.32, 2.95),
        ("image54.png", 10.25, 1.28, 1.32, 2.95),
        ("image58.png", 11.78, 1.28, 1.32, 2.95),
    ]
    for fn, x, y, w, h in pics:
        p = image_path(fn)
        if p:
            add_box(s, x - 0.06, y - 0.06, w + 0.12, h + 0.12, COL["paper"] if h < 2.5 else COL["dark"], "E1E8F0")
            add_fit_picture(s, p, x, y, w, h)
    add_bullets(s, 1.0, 4.75, [
        "Backend microservices: Identity, Catalog, Core, Trade Logistics, Blockchain.",
        "Web nghiệp vụ cho admin, supplier, manufacturer.",
        "Mobile: scan truy xuất, đánh giá, retailer order, transporter delivery.",
        "Blockchain: ghi hash, verify dữ liệu, thống kê gas; Kafka/WebSocket cho xử lý bất đồng bộ.",
    ], COL["green"])

    # Slide 12
    s = prs.slides.add_slide(blank)
    add_header(s, "Đánh giá và kết luận", "Kết thúc ngắn, để dành thời gian cho phần hỏi đáp", 12)
    add_card(s, 0.85, 1.55, 3.75, 2.0, "Ưu điểm", "Có tính thực tiễn, nhiều vai trò, luồng nghiệp vụ tương đối đầy đủ.\nBlockchain dùng hợp lý để xác minh dữ liệu.", COL["green"], COL["green_soft"], "shield")
    add_card(s, 4.85, 1.55, 3.75, 2.0, "Hạn chế", "Môi trường còn mang tính demo local.\nKiểm thử tự động và bảo mật cấu hình cần hoàn thiện thêm.", COL["orange"], COL["orange_soft"], "warning")
    add_card(s, 8.85, 1.55, 3.75, 2.0, "Hướng phát triển", "Triển khai cloud, quản lý secret chuẩn hơn, mở rộng dashboard phân tích và test tích hợp.", COL["blue"], COL["blue_soft"], "chain")
    add_box(s, 1.0, 4.62, 11.35, 0.95, COL["paper"], "E1E8F0")
    add_text(s, 1.35, 4.9, 10.7, 0.35, "Hệ thống đáp ứng mục tiêu truy xuất nguồn gốc, minh bạch dữ liệu chuỗi cung ứng và hỗ trợ chống giả sản phẩm.", 20, COL["ink"], True, "center")
    add_text(s, 1.0, 6.25, 11.35, 0.42, "Em xin chân thành cảm ơn!", 22, COL["blue"], True, "center")

    prs.save(OUT)


if __name__ == "__main__":
    make_deck()
    print(OUT)
